"""
CLI-точка входа для сборки валютного отчёта.

Usage:
    python run.py                                 # дата = предыдущий торговый день
    python run.py --date 2026-04-16               # дата задана явно
    python run.py --date 2026-04-16 --source /tmp/bonds_mock.csv
"""
import argparse
import sys
from datetime import date, datetime, timezone, timedelta
from pathlib import Path

from pptx import Presentation

sys.path.insert(0, str(Path(__file__).resolve().parent))

# Отчёт-специфичные настройки остаются в репозитории
from src.settings import (
    TEMPLATE_PATH,
    OUTPUT_DIR,
    CONFIG_DIR,
    REPORT_TITLE_LINES,
    COVER_DISCLAIMER,
)

# Общий код вынесен в пакет bonds_report_core
from bonds_report_core.sources import BONDS_DATA_URL
from bonds_report_core.config_loader import load_sections
from bonds_report_core.data_loader import load_bonds
from bonds_report_core.pdf_export import convert_pptx_to_pdf, LibreOfficeNotFoundError
from bonds_report_core.trading_day import previous_trading_day
from bonds_report_core.highlighted import load_highlighted_isins
from bonds_report_core.utils import remove_all_slides
from bonds_report_core.slides.cover import render_cover
from bonds_report_core.slides.glossary import render_glossary
from bonds_report_core.slides.disclaimer import render_disclaimer

# USD-специфичные модули и слайды
from src.slides.bond_table_currency import render_currency_section
from src.slides.yield_curve_slide import render_yield_curve
from src.moex_depth import enrich_bonds


# Путь к итоговому PDF.
# Имя фиксированное — PDF перезаписывается при каждом запуске,
# чтобы upstream-сервисы ссылались на один и тот же путь.
PDF_OUTPUT_PATH = Path("/opt/presentation_bot/out/bonds_usd.pdf")

# Имя pptx-файла. Тоже фиксированное.
PPTX_FILENAME = "bonds_usd.pptx"


def parse_args():
    p = argparse.ArgumentParser(description="Сборка валютного отчёта по облигациям.")
    p.add_argument(
        "--date",
        type=lambda s: datetime.strptime(s, "%Y-%m-%d").date(),
        default=None,
        help="Дата отчёта в формате YYYY-MM-DD. По умолчанию — текущая дата и время.",
    )
    p.add_argument(
        "--source",
        default=BONDS_DATA_URL,
        help=f"Источник CSV с данными: URL или путь к файлу. По умолчанию — {BONDS_DATA_URL}",
    )
    p.add_argument(
        "--output",
        type=Path,
        default=None,
        help=f"Путь к выходному .pptx. По умолчанию — output/{PPTX_FILENAME}.",
    )
    p.add_argument(
        "--pdf-output",
        type=Path,
        default=PDF_OUTPUT_PATH,
        help=f"Путь к выходному .pdf. По умолчанию — {PDF_OUTPUT_PATH}.",
    )
    p.add_argument(
        "--skip-pdf",
        action="store_true",
        help="Не конвертировать в PDF (для отладки).",
    )
    p.add_argument(
        "--no-depth",
        action="store_true",
        help="Не запрашивать глубину стакана из ISS (быстрый режим без Bid/Offer est.).",
    )
    return p.parse_args()


def build_presentation(report_date, source: str, output_path: Path, fetch_depth: bool = True) -> Path:
    """Основной pipeline сборки валютного отчёта."""
    if not TEMPLATE_PATH.exists():
        raise FileNotFoundError(f"Не найден шаблон: {TEMPLATE_PATH}")

    # Cover gets the full datetime (to show time); data logic uses plain date.
    date_only = report_date.date() if isinstance(report_date, datetime) else report_date

    sections = load_sections(CONFIG_DIR / "bonds.yaml")
    sections_by_id = {s.id: s for s in sections}

    # Общий список ISIN для выделения (см. рублёвый run.py).
    highlighted = load_highlighted_isins()
    if highlighted:
        print(f"Highlighted ISINs from перечня: {len(highlighted)}")

    prs = Presentation(str(TEMPLATE_PATH))
    remove_all_slides(prs)

    # --- Слайд 1: обложка ---
    render_cover(prs, report_date, REPORT_TITLE_LINES, COVER_DISCLAIMER)

    # USD- и CNY-бумаги запоминаем отдельно — они нужны для слайдов с кривыми
    # доходности, которые рисуем после всех табличных слайдов.
    usd_bonds = []
    cny_bonds = []

    # --- Валютные секции: USD, EUR, CHF, CNY ---
    # Порядок фиксирован — как в образце Газпромбанк Private.
    for sec_id in ("bonds_usd", "bonds_eur", "bonds_chf", "bonds_cny"):
        if sec_id not in sections_by_id:
            continue
        sec = sections_by_id[sec_id]
        print(f"Loading {len(sec.isins)} {sec.currency or sec_id} bonds from {source}")
        bonds = load_bonds(sec.isins, source=source)
        if not bonds:
            print(f"  [warn] ни одной бумаги {sec.currency} не нашлось в CSV — слайд пропускаю")
            continue
        if fetch_depth:
            print(f"  Fetching ISS depth for {len(bonds)} bonds...")
            enrich_bonds(bonds)
        render_currency_section(
            prs,
            bonds=bonds,
            currency=sec.currency or "",
            section_title=sec.title,
            call_dates=sec.call_dates(),
            report_date=date_only,
            highlighted=highlighted,
        )
        if sec_id == "bonds_usd":
            usd_bonds = bonds
        elif sec_id == "bonds_cny":
            cny_bonds = bonds

    # --- Кривые доходности USD и CNY ---
    # Идут после всех таблиц, перед глоссарием.
    # USD сначала, CNY — следом. Если бумаг нет — соответствующий слайд
    # не создаётся.
    if usd_bonds:
        print(f"Rendering yield curve for {len(usd_bonds)} USD bonds")
        render_yield_curve(prs, usd_bonds, date_only, currency="USD")

    if cny_bonds:
        print(f"Rendering yield curve for {len(cny_bonds)} CNY bonds")
        render_yield_curve(prs, cny_bonds, date_only, currency="CNY")

    # --- Глоссарий и дисклеймер — одинаковые с рублёвым отчётом ---
    render_glossary(prs)
    render_disclaimer(prs)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def main():
    args = parse_args()

    if args.date is None:
        args.date = datetime.now(tz=timezone(timedelta(hours=3)))

    if args.output is None:
        args.output = OUTPUT_DIR / PPTX_FILENAME

    print(f"Building USD report for {args.date} <- {args.source}")
    result = build_presentation(args.date, args.source, args.output, fetch_depth=not args.no_depth)
    print(f"Saved: {result}")

    if args.skip_pdf:
        print("[pdf] Пропускаю конвертацию (--skip-pdf)")
        return

    try:
        pdf_path = convert_pptx_to_pdf(args.output, args.pdf_output)
        print(f"PDF saved: {pdf_path}")
    except LibreOfficeNotFoundError as e:
        print(f"[pdf] Предупреждение: {e}")
    except Exception as e:
        print(f"[pdf] Ошибка конвертации: {e}")
        print(f"[pdf] pptx при этом сохранён: {args.output}")


if __name__ == "__main__":
    main()
