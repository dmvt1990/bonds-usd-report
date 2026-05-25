"""
Слайд «Важная информация и контакты» (стр. 14 образца).

Содержит:
- блок контактов сверху справа (телефоны, email)
- большой блок юридического дисклеймера на весь слайд

Контент статический, не зависит от данных.
"""
from pptx.enum.text import PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.util import Cm, Pt

from ..constants import (
    FONT_PRIMARY,
    COLOR_TEXT_PRIMARY, COLOR_TEXT_SECONDARY, COLOR_TEXT_MUTED,
    TABLE_SLIDE_TITLE_FONT_SIZE,
)
from ..fonts import set_run_font


_LAYOUT_INDEX = 3

SLIDE_TITLE = "Важная информация и контакты"

CONTACTS_LINES = [
    "8 (495) 719-19-00, Москва",
    "8 (800) 719-19-00, Россия",
    "privatebanking@gazprombank.ru",
]

# Дисклеймер состоит из абзацев. Первая строка-флажок — общий ввод.
# Дальше — пункты с маркером-кружком ◌.
DISCLAIMER_INTRO = "Данная информационная презентация (далее – «Презентация»):"

DISCLAIMER_BULLETS = [
    (
        "не является индивидуальной инвестиционной рекомендацией, и финансовые "
        "инструменты либо сделки с ними, упомянутые в данной Презентации, могут "
        "не соответствовать вашему инвестиционному профилю. Соответствие финансового "
        "инструмента либо операции инвестиционным целям, инвестиционному горизонту "
        "и отношению к риску определяется инвестором самостоятельно. Банк ГПБ (АО) "
        "не несёт ответственности за возможные убытки инвестора в случае совершения "
        "сделок с финансовыми инструментами, упомянутыми в данном документе."
    ),
    (
        "предназначена исключительно для целей информирования и не является "
        "предложением приобрести у Банка ГПБ (АО) какие-либо продукты или услуги "
        "или заключить с Банком ГПБ (АО) какие-либо сделки. Инвестиционные продукты "
        "не являются банковскими вкладами, застрахованными в рамках государственной "
        "программы страхования вкладов. В зависимости от конъюнктуры финансовых "
        "рынков стоимость инвестиционных продуктов может увеличиваться или уменьшаться. "
        "Инвестор принимает риск потери всех вложенных средств. Доходность от вложений "
        "в инвестиционные продукты за прошлые периоды не является гарантией получения "
        "такой же доходности в будущем. Перед заключением любых сделок, связанных с "
        "инвестиционными продуктами, инвестор должен оценить экономические риски и "
        "выгоды, а также юридические и налоговые последствия таких сделок, не полагаясь "
        "на мнение сотрудников Банка ГПБ (АО) и/или его дочерних и аффилированных структур."
    ),
    (
        "подготовлена исключительно в информационных целях и не является инвестиционной, "
        "финансовой, юридической, налоговой или любой иной рекомендацией, а также не "
        "является предложением покупки/продажи или подписки на представленные в ней "
        "услуги и не является основанием для заключения какого-либо договора, совершения "
        "какой-либо сделки или возникновения какого-либо обязательства или совершения "
        "юридического действия любого характера."
    ),
    (
        "не подлежит ни полному, ни частичному воспроизведению, сохранению в какой-либо "
        "информационно-поисковой системе или передаче в любой форме и любыми средствами: "
        "электронными, механическими, путём фотокопирования, записи или любым другим "
        "способом без предварительного письменного разрешения Банка ГПБ (АО)."
    ),
]

DISCLAIMER_OUTRO = (
    "Информация, представленная в Презентации, не является рекламной, как это "
    "определено в Федеральном Законе №38-ФЗ от 13 марта 2006 года «О рекламе». "
    "Презентация включает в себя информацию и анализ, не имеющие своей основной "
    "целью продвижение продукта, и не предназначена для неопределённого круга лиц. "
    "Презентация не даёт оценок в рамках Федерального Закона №135-ФЗ от 29 июля "
    "1998 года «Об оценочной деятельности в Российской Федерации». Если иное не "
    "предусмотрено действующим законодательством, закон №2300-1 от 7 февраля 1992 "
    "года «О защите прав потребителей» не применяется к настоящему документу и "
    "любым отношениям, с ним связанным."
)


# --- Геометрия ---
# Блок контактов — сверху справа
CONTACTS_LEFT = Cm(22.00)
CONTACTS_TOP = Cm(2.20)
CONTACTS_WIDTH = Cm(10.70)
CONTACTS_HEIGHT = Cm(2.00)
CONTACTS_FONT_SIZE = Pt(11)

# Блок дисклеймера — на всю ширину, ниже заголовка слайда (заголовок заканчивается около 5.0 см)
DISCLAIMER_LEFT = Cm(1.13)
DISCLAIMER_TOP = Cm(5.40)
DISCLAIMER_WIDTH = Cm(31.60)
DISCLAIMER_HEIGHT = Cm(13.00)
DISCLAIMER_FONT_SIZE = Pt(8)

BULLET_CHAR = "◌"  # та же иконка что в образце


def _fill_title(slide, text: str) -> None:
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = text
            set_run_font(run, FONT_PRIMARY)
            run.font.size = TABLE_SLIDE_TITLE_FONT_SIZE
            run.font.color.rgb = COLOR_TEXT_PRIMARY
            return


def _blank_other_placeholders(slide) -> None:
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            continue
        try:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = ""
        except Exception:
            pass


def _add_contacts_block(slide) -> None:
    """Блок с телефонами и email сверху справа."""
    box = slide.shapes.add_textbox(
        CONTACTS_LEFT, CONTACTS_TOP, CONTACTS_WIDTH, CONTACTS_HEIGHT
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    for i, line in enumerate(CONTACTS_LINES):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = line
        set_run_font(run, FONT_PRIMARY)
        run.font.size = CONTACTS_FONT_SIZE
        run.font.color.rgb = COLOR_TEXT_PRIMARY


def _add_disclaimer_block(slide) -> None:
    """Большой блок юридического текста на всю ширину."""
    box = slide.shapes.add_textbox(
        DISCLAIMER_LEFT, DISCLAIMER_TOP, DISCLAIMER_WIDTH, DISCLAIMER_HEIGHT
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    def _add_paragraph(text: str, *, bullet: bool = False, spacing_after: Pt = None):
        """Добавляет параграф с нужным форматированием."""
        if tf.paragraphs[0].text == "" and len(tf.paragraphs) == 1 and not hasattr(_add_paragraph, "_used"):
            p = tf.paragraphs[0]
            _add_paragraph._used = True
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        # Небольшой отступ после параграфа для разделения пунктов
        if spacing_after is not None:
            p.space_after = spacing_after

        run = p.add_run()
        if bullet:
            run.text = f"{BULLET_CHAR}  {text}"
        else:
            run.text = text
        set_run_font(run, FONT_PRIMARY)
        run.font.size = DISCLAIMER_FONT_SIZE
        run.font.color.rgb = COLOR_TEXT_SECONDARY

    # Вводное предложение
    _add_paragraph(DISCLAIMER_INTRO, spacing_after=Pt(4))

    # Маркированные пункты
    for bullet_text in DISCLAIMER_BULLETS:
        _add_paragraph(bullet_text, bullet=True, spacing_after=Pt(4))

    # Финальный абзац
    _add_paragraph(DISCLAIMER_OUTRO)


def render_disclaimer(prs: PresentationType) -> None:
    """Добавляет слайд 'Важная информация и контакты' в презентацию."""
    layout = prs.slide_layouts[_LAYOUT_INDEX]
    slide = prs.slides.add_slide(layout)

    _fill_title(slide, SLIDE_TITLE)
    _blank_other_placeholders(slide)

    _add_contacts_block(slide)
    _add_disclaimer_block(slide)
