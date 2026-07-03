"""
Слайд «Облигации номинированные в <валюта>» (USD / EUR / CHF / CNY).

Отличия от рублёвого отчёта:
- 13 колонок, в т.ч. колонка "Валюта расчётов" (RUB / CNY/RUB) и отдельно
  "Дата Call опциона" + "Доходность к Call опциону".
- Лот показывается как номинал в валюте актива (а не объём выпуска в млн ₽).
- Все бумаги — фиксы (флоатеров среди валютных нет в образце).
- Есть бессрочные бумаги — matdate=None → выводим "Бессрочный".

Пагинация: та же логика что у рублёвых таблиц — по 24 строки на слайд.
"""
from typing import Dict, List, Optional

from pptx.enum.text import PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.util import Cm, Pt

from ..constants import (
    FONT_PRIMARY,
    COLOR_TEXT_PRIMARY, COLOR_TEXT_SECONDARY, COLOR_TEXT_MUTED,
    TABLE_SLIDE_TITLE_FONT_SIZE,
    TABLE_SLIDE_LEFT, TABLE_SLIDE_TOP,
    TABLE_SLIDE_SOURCE_TOP, TABLE_SLIDE_SOURCE_FONT_SIZE,
    TABLE_SLIDE_FOOTNOTE_TOP, TABLE_SLIDE_FOOTNOTE_FONT_SIZE,
    TABLE_CURRENCY_COL_WIDTHS_CM,
)
from bonds_report_core.fonts import set_run_font
from bonds_report_core.formatting import (
    format_date_dmy, format_percent, format_number,
    format_integer, format_duration,
)
from bonds_report_core.models import Bond
from bonds_report_core.tables import TableColumn, render_bond_table
from ..ytm import compute_ytm


_LAYOUT_INDEX = 3  # базовый слайд с заголовком/номером страницы

# Заголовок слайда — одинаковый для всех валютных слайдов.
# Имя секции ("Облигации, номинированные в USD") идёт в шапке таблицы.
SLIDE_TITLE = "Облигации, номинированные в валюте"

CURRENCY_ROWS_PER_PAGE = 24


def _format_matdate(b: Bond) -> str:
    """Дата погашения: обычная дата или 'Бессрочный' для перпетуалов."""
    if b.matdate is None:
        return "Бессрочный"
    return format_date_dmy(b.matdate)


def _format_duration(b: Bond) -> str:
    """Дюрация: число или 'Бессрочный'."""
    if b.matdate is None:
        return "Бессрочный"
    return format_duration(b.duration_years)


def _format_lot_in_currency(b: Bond) -> str:
    """
    Минимальный лот в валюте актива.

    В MOEX ISS для валютных бумаг facevalue обычно совпадает с
    минимальным номиналом: 100, 1 000, 5 000, 100 000 или 200 000
    в валюте выпуска. Используем его (lotvalue в MOEX часто в рублях
    по курсу — бесполезен для валютных).
    """
    # Для валютных берём номинал (он уже в валюте выпуска)
    if b.facevalue is None:
        return "—"
    return format_integer(b.facevalue)


def _format_depth_est(v) -> str:
    if v is None:
        return "—"
    if v >= 1_000_000:
        return f"{v / 1_000_000:.1f}M"
    if v >= 1_000:
        return f"{v / 1_000:.0f}K"
    return f"{v:.0f}"


def _bond_to_row(i: int, b: Bond, call_date_override: Optional[str], fallback_currency: str, report_date, highlighted: set) -> List[str]:
    """
    Одна облигация → строка для таблицы валютных бумаг.
    """
    # Дата Call — из yaml или из CSV
    if call_date_override:
        call_str = call_date_override
    elif b.offerdate is not None:
        call_str = format_date_dmy(b.offerdate)
    else:
        call_str = "—"

    settle_currency = "RUB"

    ytm_value = compute_ytm(
        price=b.close,
        coupon_percent=b.couponpercent,
        matdate=b.matdate,
        report_date=report_date,
        coupon_frequency=b.couponfrequency,
    )

    mark = "✓" if b.isin in highlighted else "—"

    return [
        str(i),
        b.shortname or "—",
        b.isin,
        _format_matdate(b),
        call_str,
        format_percent(b.couponpercent, decimals=2),
        format_number(b.close, decimals=2),
        format_percent(ytm_value, decimals=2),
        _format_duration(b),
        settle_currency,
        _format_lot_in_currency(b),
        str(b.couponfrequency) if b.couponfrequency else "—",
        _format_depth_est(b.bid_est_usd),
        _format_depth_est(b.offer_est_usd),
        mark,
    ]


def _make_columns() -> List[TableColumn]:
    """Конфигурация 13 колонок таблицы валютных бумаг."""
    w = TABLE_CURRENCY_COL_WIDTHS_CM
    return [
        TableColumn("№", w[0], align="center"),
        TableColumn("Выпуск", w[1], align="center"),
        TableColumn("ISIN", w[2], align="center"),
        TableColumn("Дата\nпогашения", w[3], align="center"),
        TableColumn("Дата Call\nопциона", w[4], align="center"),
        TableColumn("Купон,\n% год.", w[5], align="center"),
        TableColumn("Цена, %\nот ном.", w[6], align="center"),
        TableColumn("Дох-ть к\nпогаш.,\n% год.", w[7], align="center"),
        TableColumn("Мод.\nдюр.", w[8], align="center"),
        TableColumn("Валюта\nрасчётов", w[9], align="center"),
        TableColumn("Мин. лот\n(валюта\nактива)", w[10], align="center"),
        TableColumn("Период.\nкупона", w[11], align="center"),
        TableColumn("Оценка\nспроса, $", w[12], align="center"),
        TableColumn("Оценка\nпредложения, $", w[13], align="center"),
        TableColumn("В\nперечне", w[14], align="center"),
    ]


def _fill_title_placeholder(slide, text: str) -> None:
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            tf = ph.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = text
            set_run_font(run, FONT_PRIMARY)
            run.font.size = TABLE_SLIDE_TITLE_FONT_SIZE
            run.font.color.rgb = COLOR_TEXT_PRIMARY
            return


def _blank_other_placeholders(slide) -> None:
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            continue
        try:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = ""
        except Exception:
            pass


def _add_source_line(slide, report_date) -> None:
    box = slide.shapes.add_textbox(
        Cm(22.00), TABLE_SLIDE_SOURCE_TOP, Cm(10.70), Cm(0.50),
    )
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"Источник: Cbonds, {format_date_dmy(report_date)}"
    set_run_font(run, FONT_PRIMARY)
    run.font.size = TABLE_SLIDE_SOURCE_FONT_SIZE
    run.font.color.rgb = COLOR_TEXT_SECONDARY


def _add_footnote(slide) -> None:
    box = slide.shapes.add_textbox(
        Cm(1.13), TABLE_SLIDE_FOOTNOTE_TOP, Cm(31.60), Cm(0.40),
    )
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "Внимание! Приведенные котировки и расчеты являются индикативными "
        "и подлежат регулярному обновлению. Отдельные параметры могут "
        "существенно изменяться под влиянием рыночной конъюнктуры."
    )
    set_run_font(run, FONT_PRIMARY)
    run.font.size = TABLE_SLIDE_FOOTNOTE_FONT_SIZE
    run.font.color.rgb = COLOR_TEXT_MUTED


def _render_one_page(
    prs: PresentationType,
    page_bonds: List[Bond],
    call_dates: Dict[str, str],
    currency: str,
    section_title: str,
    start_number: int,
    report_date,
    highlighted: set,
) -> None:
    """Один слайд-страница таблицы валютных бумаг."""
    layout = prs.slide_layouts[_LAYOUT_INDEX]
    slide = prs.slides.add_slide(layout)

    _fill_title_placeholder(slide, SLIDE_TITLE)
    _blank_other_placeholders(slide)

    columns = _make_columns()
    rows = [
        _bond_to_row(start_number + i, b, call_dates.get(b.isin), currency, report_date, highlighted)
        for i, b in enumerate(page_bonds)
    ]

    render_bond_table(
        slide,
        left=TABLE_SLIDE_LEFT,
        top=TABLE_SLIDE_TOP,
        columns=columns,
        rows=rows,
        section_title=section_title,
    )

    if report_date is not None:
        _add_source_line(slide, report_date)
    _add_footnote(slide)


def render_currency_section(
    prs: PresentationType,
    bonds: List[Bond],
    currency: str,
    section_title: str,
    call_dates: Optional[Dict[str, str]] = None,
    report_date=None,
    rows_per_page: int = CURRENCY_ROWS_PER_PAGE,
    highlighted: Optional[set] = None,
) -> None:
    """
    Добавляет серию слайдов с таблицей валютных облигаций.

    Args:
        prs: презентация
        bonds: список облигаций
        currency: 'USD' / 'EUR' / 'CHF' / 'CNY'
        section_title: "Облигации, номинированные в USD"
        call_dates: словарь ISIN → дата call-опциона (из yaml)
        report_date: дата отчёта
        rows_per_page: сколько бумаг на одном слайде
        highlighted: множество ISIN, у которых в "В перечне" стоит ✓
    """
    call_dates = call_dates or {}
    highlighted = highlighted or set()

    from datetime import date as _date
    sorted_bonds = sorted(
        bonds,
        key=lambda b: (b.matdate or _date.max),
    )

    total = len(sorted_bonds)
    if total == 0:
        return

    for page_idx in range(0, total, rows_per_page):
        page_bonds = sorted_bonds[page_idx : page_idx + rows_per_page]
        _render_one_page(
            prs,
            page_bonds=page_bonds,
            call_dates=call_dates,
            currency=currency,
            section_title=section_title,
            start_number=page_idx + 1,
            report_date=report_date,
            highlighted=highlighted,
        )
