"""
Слайды глоссария (стр. 12 и 13 образца Газпромбанк Private).

Слайд 1 глоссария: таблица с терминами и определениями (дюрация,
налоговая база по НДФЛ).
Слайд 2 глоссария: один абзац про доходность при досрочной продаже.

Контент полностью статический — не зависит ни от CSV, ни от yaml.
Чтобы изменить формулировки — правим константы в начале этого модуля.
"""
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.presentation import Presentation as PresentationType
from pptx.util import Cm, Pt

from ..constants import (
    FONT_PRIMARY,
    COLOR_TEXT_PRIMARY, COLOR_TEXT_SECONDARY,
    COLOR_ACCENT_OLIVE,
    TABLE_SLIDE_TITLE_FONT_SIZE,
)
from ..fonts import set_run_font


_LAYOUT_INDEX = 3  # тот же layout, что у табличных слайдов

SLIDE_TITLE = "Глоссарий"

# Позиция заголовка — как в placeholder'е layout (наследуется)

# --- Страница 1 глоссария ---
# Пары (термин, определение) — в образце таблица на полную ширину
GLOSSARY_PAGE1_ITEMS = [
    (
        "ДЮРАЦИЯ",
        "Средний срок до погашения денежного потока, создаваемого облигацией. "
        "Мера процентного риска. Показывает чувствительность цены облигации "
        "к изменению её доходности (на сколько процентов вырастет/упадёт "
        "цена при понижении/повышении доходности на 1%)."
    ),
    (
        "НАЛОГОВАЯ БАЗА ПО НДФЛ\nПО ОПЕРАЦИЯМ С\nНЕОБРАЩАЮЩИМИСЯ\nНА БИРЖЕ ЦЕННЫМИ БУМАГАМИ\n(В ТОМ ЧИСЛЕ ЕВРООБЛИГАЦИЯМ)",
        "Определяется отдельно от налоговой базы по обращающимся ценным бумагам. "
        "Убытки, полученные по таким операциям, не подлежат переносу на будущее.\n\n"
        "Для целей налогообложения доходы и расходы по ценным бумагам, номинированным "
        "в иностранной валюте, пересчитываются в рубли по курсу ЦБ РФ на дату получения "
        "указанных доходов (осуществления расходов). То есть положительная валютная "
        "переоценка увеличивает налогооблагаемую базу, отрицательная – уменьшает.\n\n"
        "Расходы по приобретению ценных бумаг уменьшают налоговую базу в том периоде, "
        "в котором данные ценные бумаги были реализованы. Если в отчётном периоде "
        "бумаги проданы не были, то полученный по ним купон облагается налогом в полном объёме."
    ),
]

# --- Страница 2 глоссария ---
GLOSSARY_PAGE2_TEXT = (
    "При принятии решения об инвестировании средств в облигации следует учитывать, "
    "что фактически реализованная доходность может изменяться, если облигация "
    "продаётся инвестором до объявленной даты погашения. Так, облигация, досрочно "
    "реализованная по рыночной цене, превышающей цену приобретения, может увеличить "
    "фактически реализованную доходность, и наоборот. При этом, эмитент обязан "
    "погасить облигацию по 100% от номинальной цены (цены размещения) исключительно "
    "на дату погашения."
)


# --- Геометрия ---
# Ширина полезной области слайда: 33.87 - 1.13*2 = 31.60 см
CONTENT_LEFT = Cm(1.13)
CONTENT_WIDTH = Cm(31.60)

# Таблица глоссария: колонка "термин" слева (8 см), колонка "определение" справа (23.6 см)
TERM_COL_WIDTH = Cm(8.00)
DEF_COL_WIDTH = Cm(23.60)

# Шрифты и размеры
TERM_FONT_SIZE = Pt(11)
DEF_FONT_SIZE = Pt(10)
PARAGRAPH_FONT_SIZE = Pt(12)

# Вертикальные отступы — ниже заголовка слайда (он заканчивается около 5.0 см)
GLOSSARY_TABLE_TOP = Cm(5.40)
PARAGRAPH_TOP = Cm(5.40)
PARAGRAPH_HEIGHT = Cm(10.00)


def _fill_title(slide, text: str) -> None:
    """Заполняет title-placeholder, убирает 'Образец'."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = text
            set_run_font(run, FONT_PRIMARY)
            run.font.size = TABLE_SLIDE_TITLE_FONT_SIZE
            run.font.color.rgb = COLOR_TEXT_PRIMARY
            return


def _blank_other_placeholders(slide) -> None:
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            continue
        try:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = ""
        except Exception:
            pass


def _render_glossary_page1(prs: PresentationType) -> None:
    """
    Слайд 1 глоссария — таблица из двух колонок:
    слева термин (оливковым жирным, ВЕРХНИМ регистром),
    справа определение (серым обычным).
    Без видимых границ, как на других слайдах.
    """
    layout = prs.slide_layouts[_LAYOUT_INDEX]
    slide = prs.slides.add_slide(layout)
    _fill_title(slide, SLIDE_TITLE)
    _blank_other_placeholders(slide)

    # Создаём таблицу: n строк × 2 колонки
    n_rows = len(GLOSSARY_PAGE1_ITEMS)
    table_shape = slide.shapes.add_table(
        rows=n_rows, cols=2,
        left=CONTENT_LEFT, top=GLOSSARY_TABLE_TOP,
        width=CONTENT_WIDTH, height=Cm(0.1),  # высоту PowerPoint пересчитает
    )
    tbl = table_shape.table
    tbl.columns[0].width = TERM_COL_WIDTH
    tbl.columns[1].width = DEF_COL_WIDTH

    for i, (term, definition) in enumerate(GLOSSARY_PAGE1_ITEMS):
        _fill_glossary_cell(
            tbl.cell(i, 0), term,
            font_size=TERM_FONT_SIZE,
            font_color=COLOR_ACCENT_OLIVE,
            bold=True,
        )
        _fill_glossary_cell(
            tbl.cell(i, 1), definition,
            font_size=DEF_FONT_SIZE,
            font_color=COLOR_TEXT_SECONDARY,
            bold=False,
        )

    # Убираем границы всех ячеек
    from pptx.oxml.ns import qn
    from lxml import etree
    for row in tbl.rows:
        for cell in row.cells:
            tcPr = cell._tc.get_or_add_tcPr()
            # убираем заливку по умолчанию — оставляем белый фон
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_TEXT_PRIMARY.__class__(0xFF, 0xFF, 0xFF)
            # убираем линии со всех сторон
            for side in ("lnL", "lnR", "lnT", "lnB"):
                for existing in tcPr.findall(qn(f"a:{side}")):
                    tcPr.remove(existing)
                ln = etree.SubElement(tcPr, qn(f"a:{side}"))
                ln.set("w", "0")
                etree.SubElement(ln, qn("a:noFill"))


def _fill_glossary_cell(cell, text, *, font_size, font_color, bold):
    """Стилизует ячейку таблицы глоссария."""
    cell.vertical_anchor = MSO_ANCHOR.TOP
    cell.margin_left = Cm(0.20)
    cell.margin_right = Cm(0.20)
    cell.margin_top = Cm(0.30)
    cell.margin_bottom = Cm(0.30)

    tf = cell.text_frame
    tf.word_wrap = True
    tf.clear()

    # Делим текст по \n на параграфы
    lines = text.split("\n")
    for i, line in enumerate(lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.LEFT
        run = para.add_run()
        run.text = line
        set_run_font(run, FONT_PRIMARY)
        run.font.size = font_size
        run.font.color.rgb = font_color
        run.font.bold = bold


def _render_glossary_page2(prs: PresentationType) -> None:
    """Слайд 2 глоссария — один абзац про доходность при досрочной продаже."""
    layout = prs.slide_layouts[_LAYOUT_INDEX]
    slide = prs.slides.add_slide(layout)
    _fill_title(slide, SLIDE_TITLE)
    _blank_other_placeholders(slide)

    box = slide.shapes.add_textbox(
        CONTENT_LEFT, PARAGRAPH_TOP,
        CONTENT_WIDTH, PARAGRAPH_HEIGHT,
    )
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = GLOSSARY_PAGE2_TEXT
    set_run_font(run, FONT_PRIMARY)
    run.font.size = PARAGRAPH_FONT_SIZE
    run.font.color.rgb = COLOR_TEXT_SECONDARY


def render_glossary(prs: PresentationType) -> None:
    """Добавляет два слайда глоссария в презентацию."""
    _render_glossary_page1(prs)
    _render_glossary_page2(prs)
