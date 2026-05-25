"""
Титульный слайд отчёта.

Использует Layout 0 ('1_Титульный слайд') из template.pptx — в нём уже
есть логотип Газпромбанк Private и декоративная графика.
Наша задача: заполнить заголовок, добавить дату и строку
"только для квалифицированных инвесторов".
"""
from datetime import date, datetime
from typing import List, Union

from pptx.presentation import Presentation as PresentationType
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR

from ..constants import (
    FONT_PRIMARY,
    COLOR_TEXT_PRIMARY,
    COLOR_TEXT_SECONDARY,
    COLOR_TEXT_MUTED,
    COLOR_LINE_SEPARATOR,
    FONT_SIZE_COVER_TITLE,
    FONT_SIZE_COVER_DATE,
    FONT_SIZE_COVER_DISCLAIMER,
    COVER_DATE_LEFT, COVER_DATE_TOP, COVER_DATE_WIDTH, COVER_DATE_HEIGHT,
    COVER_SEPARATOR_LEFT, COVER_SEPARATOR_TOP,
    COVER_SEPARATOR_WIDTH, COVER_SEPARATOR_HEIGHT,
    COVER_DISCLAIMER_LEFT, COVER_DISCLAIMER_TOP,
    COVER_DISCLAIMER_WIDTH, COVER_DISCLAIMER_HEIGHT,
)
from ..formatting import format_russian_date
from ..settings import REPORT_TITLE_LINES, COVER_DISCLAIMER
from ..fonts import set_run_font


# Индекс layout'а "1_Титульный слайд" в шаблоне
_COVER_LAYOUT_INDEX = 0


def _set_title_placeholder(slide, title_lines: List[str]) -> None:
    """
    Находит placeholder заголовка на слайде и заполняет его текстом.

    В layout 0 это shape с именем "Заголовок 8" (placeholder с idx=0).
    На добавленном слайде он наследуется как shape с тем же типом.
    """
    # Ищем placeholder заголовка. На копии layout'а у него обычно
    # placeholder_format.idx == 0 и has_text_frame == True
    title_shape = None
    for shape in slide.placeholders:
        # idx == 0 — это заголовок (title placeholder)
        if shape.placeholder_format.idx == 0:
            title_shape = shape
            break

    if title_shape is None:
        raise RuntimeError(
            "Не найден placeholder заголовка на титульном слайде. "
            "Проверь, что шаблон не повреждён."
        )

    tf = title_shape.text_frame
    tf.clear()
    tf.word_wrap = True

    # Первая строка — в уже существующем параграфе; остальные добавляем.
    for i, line in enumerate(title_lines):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = para.add_run()
        run.text = line
        set_run_font(run, FONT_PRIMARY)
        run.font.size = FONT_SIZE_COVER_TITLE
        run.font.color.rgb = COLOR_TEXT_PRIMARY
        run.font.bold = False


def _add_date_text(slide, report_date: Union[date, datetime]) -> None:
    """Добавляет строку с датой и временем сборки отчёта под заголовком."""
    if isinstance(report_date, datetime):
        date_str = f"{format_russian_date(report_date.date())}, {report_date.strftime('%H:%M')}"
    else:
        date_str = format_russian_date(report_date)

    box = slide.shapes.add_textbox(
        COVER_DATE_LEFT, COVER_DATE_TOP,
        COVER_DATE_WIDTH, COVER_DATE_HEIGHT,
    )
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = date_str
    set_run_font(run, FONT_PRIMARY)
    run.font.size = FONT_SIZE_COVER_DATE
    run.font.color.rgb = COLOR_TEXT_SECONDARY


def _add_separator_line(slide) -> None:
    """Тонкая горизонтальная линия-разделитель между датой и дисклеймером."""
    line = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR.STRAIGHT
        COVER_SEPARATOR_LEFT,
        COVER_SEPARATOR_TOP,
        COVER_SEPARATOR_LEFT + COVER_SEPARATOR_WIDTH,
        COVER_SEPARATOR_TOP,
    )
    line.line.color.rgb = COLOR_LINE_SEPARATOR
    line.line.width = Pt(0.5)


def _add_disclaimer_text(slide) -> None:
    """Мелкая подпись 'только для квалифицированных инвесторов'."""
    box = slide.shapes.add_textbox(
        COVER_DISCLAIMER_LEFT, COVER_DISCLAIMER_TOP,
        COVER_DISCLAIMER_WIDTH, COVER_DISCLAIMER_HEIGHT,
    )
    tf = box.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = COVER_DISCLAIMER
    set_run_font(run, FONT_PRIMARY)
    run.font.size = FONT_SIZE_COVER_DISCLAIMER
    run.font.color.rgb = COLOR_TEXT_MUTED


def render_cover(prs: PresentationType, report_date: Union[date, datetime]) -> None:
    """
    Добавляет титульный слайд в презентацию.

    Args:
        prs: открытая презентация (на основе template.pptx)
        report_date: дата (или datetime) сборки отчёта. datetime → '25 мая 2026 г., 15:30'
    """
    layout = prs.slide_layouts[_COVER_LAYOUT_INDEX]
    slide = prs.slides.add_slide(layout)

    _set_title_placeholder(slide, REPORT_TITLE_LINES)
    _add_date_text(slide, report_date)
    _add_separator_line(slide)
    _add_disclaimer_text(slide)
