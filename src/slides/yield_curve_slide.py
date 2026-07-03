"""
Слайд «Кривая доходности USD» — новая страница в bonds_usd.pptx.

На слайде — крупная картинка графика (через yield_curve.py),
заголовок в стиле других слайдов, и стандартный footnote снизу.
"""
from datetime import date
from typing import List

from pptx.enum.text import PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.util import Cm

from ..constants import (
    FONT_PRIMARY,
    COLOR_TEXT_PRIMARY, COLOR_TEXT_SECONDARY, COLOR_TEXT_MUTED,
    TABLE_SLIDE_TITLE_FONT_SIZE,
    TABLE_SLIDE_SOURCE_TOP, TABLE_SLIDE_SOURCE_FONT_SIZE,
    TABLE_SLIDE_FOOTNOTE_TOP, TABLE_SLIDE_FOOTNOTE_FONT_SIZE,
)
from bonds_report_core.fonts import set_run_font
from bonds_report_core.formatting import format_date_dmy
from bonds_report_core.models import Bond
from ..yield_curve import build_yield_curve_chart


_LAYOUT_INDEX = 3  # тот же layout, что у табличных слайдов

# Геометрия картинки графика на слайде.
# Слайд — 33.87 × 19.05 см. Заголовок сверху до ~5.4 см,
# подпись «Источник» на 17.2 см, footnote на 17.9 см.
# График должен заканчиваться не ниже 16.6 см — даём 11.5 см высоты.
CHART_LEFT = Cm(2.70)
CHART_TOP = Cm(5.10)
CHART_WIDTH = Cm(28.50)
CHART_HEIGHT = Cm(11.50)


def _fill_title(slide, text: str) -> None:
    """Заполняет title-placeholder."""
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            tf = ph.text_frame
            tf.clear()
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = text
            set_run_font(run, FONT_PRIMARY)
            run.font.size = TABLE_SLIDE_TITLE_FONT_SIZE
            run.font.color.rgb = COLOR_TEXT_PRIMARY
            return


def _blank_other_placeholders(slide) -> None:
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            continue
        try:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = ""
        except Exception:
            pass


def _add_source_line(slide, report_date) -> None:
    """Подпись 'Источник: Cbonds, DD.MM.YYYY' в правом нижнем углу."""
    box = slide.shapes.add_textbox(
        Cm(22.00), TABLE_SLIDE_SOURCE_TOP, Cm(10.70), Cm(0.50),
    )
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"Источник: Cbonds, {format_date_dmy(report_date)}"
    set_run_font(run, FONT_PRIMARY)
    run.font.size = TABLE_SLIDE_SOURCE_FONT_SIZE
    run.font.color.rgb = COLOR_TEXT_SECONDARY


def _add_footnote(slide) -> None:
    """Стандартный disclaimer внизу слайда."""
    box = slide.shapes.add_textbox(
        Cm(1.13), TABLE_SLIDE_FOOTNOTE_TOP, Cm(31.60), Cm(0.40),
    )
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = (
        "Внимание! Приведенные котировки и расчеты являются индикативными "
        "и подлежат регулярному обновлению. Отдельные параметры могут "
        "существенно изменяться под влиянием рыночной конъюнктуры."
    )
    set_run_font(run, FONT_PRIMARY)
    run.font.size = TABLE_SLIDE_FOOTNOTE_FONT_SIZE
    run.font.color.rgb = COLOR_TEXT_MUTED


def render_yield_curve(
    prs: PresentationType,
    bonds: List[Bond],
    report_date: date,
    currency: str = "USD",
) -> None:
    """
    Добавляет слайд с кривой доходности в презентацию.

    Если бумаг недостаточно для рисования графика — слайд не создаётся,
    функция просто ничего не делает. Это удобно если в отчёте
    только пара бумаг или данные CSV частично сломаны.

    Args:
        prs: презентация на основе template.pptx
        bonds: список облигаций одной валюты (отфильтрованных перед вызовом)
        report_date: дата отчёта
        currency: код валюты для заголовка ('USD', 'CNY', ...)
    """
    png_buf = build_yield_curve_chart(bonds, report_date, currency=currency)
    if png_buf is None:
        return

    layout = prs.slide_layouts[_LAYOUT_INDEX]
    slide = prs.slides.add_slide(layout)

    _fill_title(slide, f"Кривая доходности {currency}")
    _blank_other_placeholders(slide)

    slide.shapes.add_picture(
        png_buf,
        CHART_LEFT, CHART_TOP,
        width=CHART_WIDTH, height=CHART_HEIGHT,
    )

    _add_source_line(slide, report_date)
    _add_footnote(slide)
